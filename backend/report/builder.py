"""保存済み結果から PPTX のスライドを組み立てる本体。

**新しい集計をしない。** 各モジュールの ``analysis.py`` が既に計算し、保存済みの
json / csv に入っている値だけを使う（並べ替えと上位 N 件の切り出しは行うが、
数値を作り直すことはしない）。

章立ては ``analysis.SECTION_ORDER``（Hang AP → Floor Peak → RRM）で固定。
選ばれなかったモジュールのスライドは作らない。

体裁は実用的な無地。HPE ブランドのテンプレートは使わない（35 番の前提 7）。

ネットワークアクセス・LLM 呼び出しは行わない。
"""
from __future__ import annotations

from datetime import datetime
from typing import Any, Sequence

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from . import charts

# ---------------------------------------------------------------------------
# 体裁
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN = Inches(0.55)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN * 2

#: 見出しの下端（本文はここから始める）
BODY_TOP = Inches(1.35)
BODY_HEIGHT = SLIDE_HEIGHT - BODY_TOP - Inches(0.5)

TEXT_COLOR = RGBColor(0x1F, 0x29, 0x37)
MUTED_COLOR = RGBColor(0x6B, 0x72, 0x80)
ACCENT_COLOR = RGBColor(0x0F, 0x62, 0x8B)
WARN_COLOR = RGBColor(0xB4, 0x54, 0x09)
HEADER_FILL = RGBColor(0xE8, 0xEE, 0xF3)

#: 明細の要約に載せる行数。全行は 1 枚に収まらないので必ず切る（切ったことは書く）
DETAIL_ROWS = 15

#: 積み上げ棒に描くバケット数の上限。超えたぶんは新しい側を残す
MAX_CHART_BUCKETS = 48

#: フロア別グラフを作るフロア数の上限（超えたら注記する）
MAX_FLOOR_CHARTS = 8

#: インパクトのスライドに載せる AP 数
IMPACT_AP_ROWS = 10


def bar_height(count: int, available: Emu, *, per=Inches(0.42), base=Inches(1.5)) -> Emu:
    """横棒グラフの高さ。項目が少ないときに棒が異様に太くならないよう縮める。

    上限は ``available``（本文に使える高さ）。項目が多いときは切り上げず、
    目盛ラベルのサイズで詰める（:func:`charts._label_size`）。
    """
    return Emu(min(int(available), int(base) + int(per) * max(count, 1)))


# ---------------------------------------------------------------------------
# 値の整形
# ---------------------------------------------------------------------------


def text_of(value: object, *, blank: str = "-") -> str:
    """セル・箇条書きに出せる文字列にする（欠測は ``blank``）。"""
    if value is None:
        return blank
    if isinstance(value, pd.Timestamp):
        return blank if pd.isna(value) else value.strftime("%Y-%m-%d %H:%M:%S")
    try:
        if pd.isna(value):
            return blank
    except (TypeError, ValueError):
        pass
    if isinstance(value, bool):
        return "はい" if value else "いいえ"
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    text = str(value).strip()
    return text or blank


def period_text(period: object) -> str:
    """``["2026-01-01 00:00:00", "2026-01-02 00:00:00"]`` を 1 行にする。"""
    if not isinstance(period, (list, tuple)) or len(period) != 2:
        return "-"
    first, last = (text_of(p) for p in period)
    if first == "-" and last == "-":
        return "-"
    return f"{first} 〜 {last}"


def counts_text(counts: object) -> str:
    """``{"回復": 3, "継続中": 1}`` を 1 行にする。"""
    if not isinstance(counts, dict) or not counts:
        return "-"
    return " / ".join(f"{k} {v}" for k, v in counts.items())


def window_text(meta: dict[str, Any]) -> str:
    start = text_of(meta.get("window_start"), blank="")
    end = text_of(meta.get("window_end"), blank="")
    if not start and not end:
        return "全データ（期間指定なし）"
    return f"{start or '（開始指定なし）'} 〜 {end or '（終了指定なし）'}"


# ---------------------------------------------------------------------------
# スライドの部品
# ---------------------------------------------------------------------------


def new_presentation() -> Presentation:
    """16:9 の空のプレゼンテーション。"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_slide(prs: Presentation):
    """白紙レイアウトのスライドを足す（既定テンプレートの 7 番目が白紙）。"""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def _set(paragraph, text: str, *, size: Pt, bold: bool = False, color: RGBColor = TEXT_COLOR):
    run = paragraph.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def add_heading(slide, section_label: str, title: str) -> None:
    """スライド上部の見出し（章名 + タイトル）と区切り線。"""
    frame = _textbox(slide, MARGIN, Inches(0.35), CONTENT_WIDTH, Inches(0.9))
    first = frame.paragraphs[0]
    _set(first, section_label, size=Pt(12), bold=True, color=ACCENT_COLOR)
    second = frame.add_paragraph()
    _set(second, title, size=Pt(22), bold=True)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.22), CONTENT_WIDTH, Emu(12700)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    line.shadow.inherit = False


def add_lines(
    slide,
    lines: Sequence[tuple[str, str]] | Sequence[str],
    *,
    left=MARGIN,
    top=BODY_TOP,
    width=CONTENT_WIDTH,
    height=Inches(2.4),
    size: Pt = Pt(13),
) -> None:
    """``(見出し, 値)`` の並び、または文字列の並びを段落として置く。"""
    frame = _textbox(slide, left, top, width, height)
    for index, item in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(4)
        if isinstance(item, tuple):
            label, value = item
            _set(paragraph, f"{label}: ", size=size, bold=True, color=MUTED_COLOR)
            _set(paragraph, value, size=size)
        else:
            _set(paragraph, str(item), size=size)


def add_note(slide, text: str, *, top=None, color: RGBColor = MUTED_COLOR) -> None:
    """スライド下端の注記（件数を切ったこと・警告など、黙って隠さないための行）。"""
    frame = _textbox(
        slide, MARGIN, top if top is not None else SLIDE_HEIGHT - Inches(0.62),
        CONTENT_WIDTH, Inches(0.45),
    )
    _set(frame.paragraphs[0], text, size=Pt(10), color=color)


def add_table(
    slide,
    header: Sequence[str],
    rows: Sequence[Sequence[object]],
    *,
    left=MARGIN,
    top=BODY_TOP,
    width=CONTENT_WIDTH,
    height=Inches(2.0),
    col_ratios: Sequence[float] | None = None,
    size: Pt = Pt(11),
):
    """見出し行つきの表。``col_ratios`` で列幅の比を指定できる。"""
    shape = slide.shapes.add_table(len(rows) + 1, len(header), left, top, width, height)
    table = shape.table

    if col_ratios:
        total = sum(col_ratios)
        for index, ratio in enumerate(col_ratios):
            table.columns[index].width = Emu(int(width * ratio / total))

    for index, label in enumerate(header):
        cell = table.cell(0, index)
        cell.text = str(label)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        paragraph = cell.text_frame.paragraphs[0]
        for run in paragraph.runs:
            run.font.size = size
            run.font.bold = True
            run.font.color.rgb = TEXT_COLOR

    for r, row in enumerate(rows, start=1):
        table.rows[r].height = Inches(0.3)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text_of(value)
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = size
                run.font.color.rgb = TEXT_COLOR
    return table


def add_color_legend(
    slide, entries: Sequence[tuple[str, object]], *, left, top, width, size: Pt = Pt(10)
) -> None:
    """色見本 + 名前の凡例（棒ごとに色を変えたグラフは標準の凡例が使えない）。"""
    swatch = Inches(0.14)
    step = Inches(0.24)
    for index, (label, color) in enumerate(entries):
        y = top + step * index
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y + Inches(0.03), swatch, swatch)
        box.fill.solid()
        box.fill.fore_color.rgb = charts._rgb(color)
        box.line.fill.background()
        box.shadow.inherit = False
        frame = _textbox(slide, left + swatch + Inches(0.08), y - Inches(0.02), width, step)
        frame.margin_top = 0
        frame.margin_bottom = 0
        _set(frame.paragraphs[0], str(label), size=size)


# ---------------------------------------------------------------------------
# 表紙
# ---------------------------------------------------------------------------


def build_cover(prs: Presentation, sources: Sequence[Any], generated_at: datetime) -> str:
    """表紙。**常に作る。** 含まれるモジュールを固定順で並べる。"""
    slide = add_slide(prs)
    frame = _textbox(slide, MARGIN, Inches(1.5), CONTENT_WIDTH, Inches(1.6))
    _set(frame.paragraphs[0], "RF モニタリング 横断レポート", size=Pt(34), bold=True)
    second = frame.add_paragraph()
    _set(
        second,
        "Hang AP / Floor Peak / RRM の保存済み分析結果をまとめたものです（再分析はしていません）",
        size=Pt(13), color=MUTED_COLOR,
    )

    add_lines(
        slide,
        [("生成日時", generated_at.strftime("%Y-%m-%d %H:%M:%S UTC"))],
        top=Inches(3.1), height=Inches(0.5),
    )

    add_table(
        slide,
        ["章", "分析結果", "保存日時", "分析条件"],
        [
            [
                source.label,
                source.name,
                text_of(source.meta.get("saved_at")),
                _condition_summary(source.meta),
            ]
            for source in sources
        ],
        top=Inches(3.7),
        height=Inches(0.4) * (len(sources) + 1),
        col_ratios=(1.4, 3.0, 2.0, 6.0),
        size=Pt(10),
    )
    return "RF モニタリング 横断レポート"


def _condition_summary(meta: dict[str, Any]) -> str:
    """表紙に 1 行で載せる分析条件（長いので先頭だけ）。"""
    text = str(meta.get("condition_text") or "").replace("\n", " / ").strip()
    return (text[:110] + "…") if len(text) > 110 else (text or "-")


# ---------------------------------------------------------------------------
# Hang AP
# ---------------------------------------------------------------------------

HANGAP_DETAIL_COLUMNS: tuple[str, ...] = (
    "ap_name", "site_name", "ゼロ開始", "ゼロ終了", "連続ゼロ回数", "回復状況", "周辺AP判定",
)
#: 上位 N 件を決める列。**降順に並べ替えるだけ**で、新しい値は作らない
HANGAP_SORT_COLUMN = "連続ゼロ回数"

HANGAP_PREV_CLIENTS_DETAIL_COLUMNS: tuple[str, ...] = (
    "ap_name", "site_name", "直前clients", "AP最大clients", "連続ゼロ回数", "回復状況", "周辺AP判定",
)
#: 上位 N 件を決める列（直前clients の多い順スライド用）
HANGAP_PREV_CLIENTS_SORT_COLUMN = "直前clients"


def build_hangap(prs: Presentation, source: Any) -> list[tuple[str, str]]:
    """Hang AP の章（サマリ + 明細の要約）。``(kind, title)`` の並びを返す。"""
    meta, rows = source.meta, source.rows
    label = source.label
    built: list[tuple[str, str]] = []

    slide = add_slide(prs)
    title = "Hang AP — サマリ"
    add_heading(slide, label, title)
    add_lines(
        slide,
        [
            ("対象期間（メトリクス）", period_text(meta.get("metrics_period"))),
            ("対象期間（イベント）", period_text(meta.get("events_period"))),
            ("検知件数（ゼロ区間）", f"{meta.get('detected_intervals', 0)} 件"),
            ("AP イベントと一致した区間", f"{meta.get('event_matched_intervals', 0)} 件"),
            ("退場疑い", f"{meta.get('exodus_suspected', 0)} 件"),
            ("対象 AP / 走査ファイル", f"{meta.get('ap_count', 0)} 台 / {meta.get('files_scanned', 0)} ファイル"),
        ],
        height=Inches(2.1),
    )
    add_table(
        slide, ["回復状況", "件数"],
        [[k, v] for k, v in (meta.get("recovery_status") or {}).items()] or [["（記録なし）", "-"]],
        top=Inches(3.6), width=Inches(5.8), height=Inches(0.34),
        col_ratios=(3.0, 1.0),
    )
    add_table(
        slide, ["周辺 AP 判定", "件数"],
        [[k, v] for k, v in (meta.get("neighbor_verdict") or {}).items()] or [["（記録なし）", "-"]],
        left=MARGIN + Inches(6.4), top=Inches(3.6), width=Inches(5.8), height=Inches(0.34),
        col_ratios=(3.0, 1.0),
    )
    add_note(slide, _warning_note(meta), color=_warning_color(meta))
    built.append(("summary", title))

    slide = add_slide(prs)
    title = f"Hang AP — 明細の要約（{HANGAP_SORT_COLUMN}の多い順 上位 {DETAIL_ROWS} 件）"
    add_heading(slide, label, title)
    top_rows = _top_rows(rows, HANGAP_SORT_COLUMN, DETAIL_ROWS)
    if top_rows.empty:
        add_lines(slide, ["（検知された区間はありません）"])
    else:
        add_table(
            slide,
            list(HANGAP_DETAIL_COLUMNS),
            [
                [row[column] for column in HANGAP_DETAIL_COLUMNS]
                for _, row in top_rows.iterrows()
            ],
            height=Inches(0.32) * (len(top_rows) + 1),
            col_ratios=(2.4, 2.4, 2.4, 2.4, 1.4, 1.6, 2.0),
            size=Pt(10),
        )
    add_note(
        slide,
        f"全 {len(rows)} 件のうち上位 {min(DETAIL_ROWS, len(rows))} 件のみ。"
        "全件は Hang AP の保存済み結果（xlsx / csv）を参照してください",
    )
    built.append(("table", title))

    slide = add_slide(prs)
    title = f"Hang AP — 明細の要約（{HANGAP_PREV_CLIENTS_SORT_COLUMN}の多い順 上位 {DETAIL_ROWS} 件）"
    add_heading(slide, label, title)
    top_rows = _top_rows(rows, HANGAP_PREV_CLIENTS_SORT_COLUMN, DETAIL_ROWS, tiebreak="ap_name")
    if top_rows.empty:
        add_lines(slide, ["（検知された区間はありません）"])
    else:
        add_table(
            slide,
            list(HANGAP_PREV_CLIENTS_DETAIL_COLUMNS),
            [
                [row[column] for column in HANGAP_PREV_CLIENTS_DETAIL_COLUMNS]
                for _, row in top_rows.iterrows()
            ],
            height=Inches(0.32) * (len(top_rows) + 1),
            col_ratios=(2.4, 2.4, 1.3, 1.3, 1.3, 1.6, 2.0),
            size=Pt(10),
        )
    add_note(
        slide,
        f"全 {len(rows)} 件のうち上位 {min(DETAIL_ROWS, len(rows))} 件のみ。"
        "全件は Hang AP の保存済み結果（xlsx / csv）を参照してください",
    )
    built.append(("table", title))
    return built


def _top_rows(
    rows: pd.DataFrame, column: str, limit: int, *, tiebreak: str | None = None,
) -> pd.DataFrame:
    """``column`` の降順で上位 ``limit`` 件。列が無ければ先頭から取る。

    ``tiebreak`` を指定すると、``column`` が同値の行は ``tiebreak`` 列の昇順で並べる。
    """
    if rows.empty:
        return rows
    if column not in rows.columns:
        return rows.head(limit)
    if tiebreak and tiebreak in rows.columns:
        return rows.sort_values(
            [column, tiebreak], ascending=[False, True], na_position="last",
        ).head(limit)
    return rows.sort_values(column, ascending=False, na_position="last").head(limit)


def _warning_note(meta: dict[str, Any]) -> str:
    warnings = list(meta.get("warnings") or [])
    if not warnings:
        return "警告なし"
    head = warnings[0].replace("\n", " ")
    if len(warnings) > 1:
        return f"警告 {len(warnings)} 件（先頭のみ）: {head}"
    return f"警告 1 件: {head}"


def _warning_color(meta: dict[str, Any]) -> RGBColor:
    return WARN_COLOR if (meta.get("warnings") or []) else MUTED_COLOR


# ---------------------------------------------------------------------------
# Floor Peak
# ---------------------------------------------------------------------------

FLOORPEAK_TOP_N = 20


def build_floorpeak(prs: Presentation, source: Any) -> list[tuple[str, str]]:
    """Floor Peak の章（サマリ + フロアごとのトップ 20 グラフ）。"""
    meta, rows = source.meta, source.rows
    label = source.label
    built: list[tuple[str, str]] = []

    slide = add_slide(prs)
    title = "Floor Peak — サマリ"
    add_heading(slide, label, title)
    add_lines(
        slide,
        [
            ("対象サイト", text_of(meta.get("site_label") or meta.get("site_name"))),
            ("ピーク時刻", text_of(meta.get("peak_time"))),
            ("ピーク時点の合計接続端末数", f"{meta.get('peak_total_clients', 0)} 台"),
            ("対象フロア数 / AP 台数", f"{meta.get('floor_count', 0)} フロア / {meta.get('ap_count', 0)} 台"),
            ("ピークの選定方法", text_of(meta.get("selected_by"))),
            ("フロアマップ", text_of(meta.get("floormap_file"))),
        ],
        height=Inches(2.1),
    )
    floors = list(meta.get("floors") or [])
    add_table(
        slide,
        ["フロア", "AP 台数", "接続端末数"],
        [[f.get("map_name"), f.get("ap_count"), f.get("num_clients")] for f in floors]
        or [["（フロアの記録がありません）", "-", "-"]],
        top=Inches(3.6),
        width=Inches(7.4),
        height=Inches(0.32) * (len(floors) + 1),
        col_ratios=(4.0, 1.6, 1.8),
    )
    add_note(slide, _warning_note(meta), color=_warning_color(meta))
    built.append(("summary", title))

    for floor in _chart_floors(meta, rows):
        built.append(_floor_chart_slide(prs, source, floor, len(floors)))
    return built


def _chart_floors(meta: dict[str, Any], rows: pd.DataFrame) -> list[str]:
    """グラフを作るフロア。**保存結果の floors の順**をそのまま使う。"""
    names = [str(f.get("map_name")) for f in (meta.get("floors") or [])]
    if not names and not rows.empty and "map_name" in rows.columns:
        names = list(dict.fromkeys(str(v) for v in rows["map_name"]))
    return names[:MAX_FLOOR_CHARTS]


def _floor_chart_slide(prs, source, floor: str, floor_total: int) -> tuple[str, str]:
    meta, rows = source.meta, source.rows
    title = f"Floor Peak — {floor} のトップ {FLOORPEAK_TOP_N}"
    slide = add_slide(prs)
    add_heading(slide, source.label, title)

    top = _floor_rows(rows, floor)
    if top.empty:
        add_lines(slide, [f"（{floor} に AP がありません）"])
        return ("chart", title)

    colors = meta.get("model_colors") or {}
    default_color = meta.get("default_model_color") or charts.FALLBACK_COLOR
    models = [str(m) for m in top["model"]]

    charts.add_bar(
        slide,
        left=MARGIN, top=BODY_TOP, width=CONTENT_WIDTH - Inches(2.1),
        height=bar_height(len(top), BODY_HEIGHT - Inches(0.3)),
        title=f"{floor} / ピーク時点 {text_of(meta.get('peak_time'))} の接続端末数",
        series_name="接続端末数",
        categories=[str(name) for name in top["ap_name"]],
        values=[int(v) for v in top["num_clients"]],
        point_colors=[colors.get(model) or default_color for model in models],
    )
    add_color_legend(
        slide,
        [(model, colors.get(model) or default_color) for model in dict.fromkeys(models)],
        left=SLIDE_WIDTH - MARGIN - Inches(1.9), top=BODY_TOP + Inches(0.1), width=Inches(1.8),
    )
    note = (
        f"棒の色は AP のモデル（色の定義は Floor Peak の分析結果に含まれるもの）。"
        f"このフロアの AP {int((rows['map_name'] == floor).sum())} 台のうち上位 {len(top)} 台"
    )
    if floor_total > MAX_FLOOR_CHARTS:
        note += f" / グラフは先頭 {MAX_FLOOR_CHARTS} フロアのみ（全 {floor_total} フロア）"
    add_note(slide, note)
    return ("chart", title)


def _floor_rows(rows: pd.DataFrame, floor: str) -> pd.DataFrame:
    """このフロアのトップ N。**保存済みの rank_in_floor をそのまま使う。**"""
    if rows.empty or "map_name" not in rows.columns:
        return rows.head(0)
    target = rows[rows["map_name"].astype(str) == str(floor)]
    if "rank_in_floor" in target.columns:
        target = target.sort_values("rank_in_floor")
    return target.head(FLOORPEAK_TOP_N)


# ---------------------------------------------------------------------------
# RRM
# ---------------------------------------------------------------------------


def build_rrm(prs: Presentation, source: Any) -> list[tuple[str, str]]:
    """RRM の章（サマリ + 時間帯別の積み上げ棒 + インパクト）。"""
    meta = source.meta
    label = source.label
    built: list[tuple[str, str]] = []

    slide = add_slide(prs)
    title = "RRM — サマリ"
    add_heading(slide, label, title)
    site_names = list(meta.get("site_labels") or meta.get("site_names") or [])
    add_lines(
        slide,
        [
            ("対象期間", window_text(meta)),
            ("対象サイト", "、".join(site_names) if site_names else "すべて"),
            ("チャネル変更", f"{meta.get('change_count', 0)} 件（{counts_text(meta.get('changes_by_class'))}）"),
            ("評価のみ（no-op）", f"{meta.get('noop_count', 0)} 件（{counts_text(meta.get('noop_by_class'))}）"),
            ("レーダー検知", f"{meta.get('radar_detected', 0)} 件"
                            f"（変更あり {meta.get('radar_with_change', 0)} / "
                            f"ACTION 未記録 {meta.get('radar_without_action', 0)}）"),
            ("照合不可 / 汚染", f"{meta.get('unmatched_count', 0)} 件 / {meta.get('contaminated_count', 0)} 件"),
            ("インパクト合計（変更前の接続端末数）", f"{meta.get('impact_total', 0)} 台"),
        ],
        height=Inches(2.4),
    )
    by_class = list(meta.get("by_classification") or [])
    add_table(
        slide,
        ["分類", "イベント", "変更", "no-op", "汚染", "照合不可", "インパクト合計", "インパクト平均"],
        [
            [
                item.get("classification"), item.get("events"), item.get("changes"),
                item.get("noop"), item.get("contaminated"), item.get("unmatched"),
                item.get("impact_total"), _round(item.get("impact_avg")),
            ]
            for item in by_class
        ] or [["（分類別の記録がありません）", "-", "-", "-", "-", "-", "-", "-"]],
        top=Inches(4.0),
        height=Inches(0.32) * (len(by_class) + 1),
        col_ratios=(2.0, 1.2, 1.2, 1.2, 1.2, 1.4, 1.8, 1.8),
    )
    add_note(slide, _warning_note(meta), color=_warning_color(meta))
    built.append(("summary", title))

    built.append(_rrm_hourly_slide(prs, source))
    built.append(_rrm_impact_slide(prs, source))
    return built


def _round(value: object) -> object:
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        return round(float(value), 2)
    return value


def bucket_label(value: object) -> str:
    """バケットの目盛ラベル。秒は落とす（横棒では左端の幅がそのまま作画領域を削る）。"""
    text = text_of(value)
    return text[:-3] if len(text) == 19 and text.endswith(":00") else text


def rrm_chart_buckets(meta: dict[str, Any]) -> list[dict[str, Any]]:
    """グラフに描くバケット。多すぎるときは **新しい側** を残す。"""
    hourly = list(meta.get("hourly") or [])
    return hourly[-MAX_CHART_BUCKETS:] if len(hourly) > MAX_CHART_BUCKETS else hourly


def rrm_classifications(meta: dict[str, Any]) -> list[str]:
    """系列の並び。**分類の定義はバックエンドの保存結果**を使う。"""
    names = list(meta.get("classifications") or [])
    if names:
        return names
    return list((meta.get("changes_by_class") or {}).keys())


def _rrm_hourly_slide(prs, source) -> tuple[str, str]:
    meta = source.meta
    title = "RRM — 時間帯別のチャネル変更回数"
    slide = add_slide(prs)
    add_heading(slide, source.label, title)

    buckets = rrm_chart_buckets(meta)
    names = rrm_classifications(meta)
    if not buckets or not names:
        add_lines(slide, ["（時間帯別の記録がありません）"])
        return ("chart", title)

    colors = meta.get("class_colors") or {}
    charts.add_stacked_bar(
        slide,
        left=MARGIN, top=BODY_TOP, width=CONTENT_WIDTH,
        height=bar_height(len(buckets), BODY_HEIGHT - Inches(0.15)),
        title="時間帯別のチャネル変更回数（分類別の積み上げ / 1 時間バケット）",
        categories=[bucket_label(b.get("bucket")) for b in buckets],
        series=[
            (name, [int(b.get(f"changes_{name}") or 0) for b in buckets])
            for name in names
        ],
        colors=[colors.get(name) or charts.FALLBACK_COLOR for name in names],
    )
    total = len(meta.get("hourly") or [])
    note = f"1 時間バケット / 全 {total} バケット"
    if total > len(buckets):
        note = f"グラフは新しい側 {len(buckets)} バケットのみ（全 {total} バケット）"
    add_note(slide, note)
    return ("chart", title)


def _rrm_impact_slide(prs, source) -> tuple[str, str]:
    meta = source.meta
    title = "RRM — インパクト（変更前の接続端末数）"
    slide = add_slide(prs)
    add_heading(slide, source.label, title)

    by_class = list(meta.get("by_classification") or [])
    colors = meta.get("class_colors") or {}
    if by_class:
        charts.add_bar(
            slide,
            left=MARGIN, top=BODY_TOP, width=Inches(6.6),
            height=bar_height(len(by_class), Inches(4.6)),
            title="分類別のインパクト合計",
            series_name="インパクト合計",
            categories=[str(item.get("classification")) for item in by_class],
            values=[int(item.get("impact_total") or 0) for item in by_class],
            point_colors=[
                colors.get(str(item.get("classification"))) or charts.FALLBACK_COLOR
                for item in by_class
            ],
        )
    else:
        add_lines(slide, ["（分類別の記録がありません）"], height=Inches(0.5))

    by_ap = list(meta.get("by_ap") or [])[:IMPACT_AP_ROWS]
    add_table(
        slide,
        ["AP", "サイト", "変更", "インパクト"],
        [
            [item.get("ap_name"), item.get("site_name"), item.get("changes"), item.get("impact_total")]
            for item in by_ap
        ] or [["（AP 別の記録がありません）", "-", "-", "-"]],
        left=MARGIN + Inches(7.0),
        top=BODY_TOP,
        width=Inches(5.2),
        height=Inches(0.3) * (len(by_ap) + 1),
        col_ratios=(2.4, 2.4, 1.2, 1.4),
        size=Pt(10),
    )
    add_note(
        slide,
        f"インパクトは変更直前の接続端末数の合計（全体 {meta.get('impact_total', 0)} 台）。"
        f"表は AP 別の上位 {len(by_ap)} 件（変更回数の多い順、保存済み結果の並びのまま）",
    )
    return ("chart", title)


#: 章 → 組み立て関数。``analysis`` はこの表を通してしか章を作らない
SECTION_BUILDERS = {
    "hangap": build_hangap,
    "floorpeak": build_floorpeak,
    "rrm": build_rrm,
}
