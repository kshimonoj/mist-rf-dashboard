"""python-pptx のネイティブグラフを組み立てるヘルパ。

**ブラウザ（recharts）の見た目をキャプチャしない。** 保存済み結果に入っている
数値から PowerPoint のグラフとして描き直す（35 番の前提 4）。

描画の確認（LibreOffice で PDF 化して目視）で分かったこと:

- python-pptx は openpyxl の ``delete=False`` にあたる落とし穴を持たない。
  軸・凡例・色分けは既定のまま正しく出る（27 番で xlsx が踏んだ軸消失は起きない）。
- ただし **目盛ラベルの既定サイズ（18pt）は大きすぎる**。カテゴリが 20 件を
  超えると LibreOffice がラベルを 1 つおきに間引く。項目数に応じてラベルの
  フォントサイズを下げること（:func:`_label_size`）。
- カテゴリ軸は既定で「先頭が下」になる。時系列も順位も上から読ませたいので、
  ``c:scaling/c:orientation val="maxMin"`` を入れて反転させる
  （python-pptx に相当のプロパティが無いため XML を直接触る）。

ネットワークアクセス・LLM 呼び出しは行わない。
"""
from __future__ import annotations

from typing import Sequence

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

#: 系列の色が取れなかったときの灰色（xlsx 側の既定と合わせる）
FALLBACK_COLOR = "9E9E9E"

#: 棒の間隔（27 番で xlsx に入れた gapWidth=40 と同じ見た目にする）
GAP_WIDTH = 40


def _rgb(color: object) -> RGBColor:
    """``RRGGBB`` / ``#RRGGBB`` を :class:`RGBColor` にする。読めなければ灰色。"""
    text = str(color or "").strip().lstrip("#")
    if len(text) != 6:
        text = FALLBACK_COLOR
    try:
        return RGBColor.from_string(text.upper())
    except ValueError:
        return RGBColor.from_string(FALLBACK_COLOR)


def _label_size(count: int) -> Pt:
    """カテゴリ数に応じた目盛ラベルのサイズ。

    既定（18pt）のままだと LibreOffice / PowerPoint がラベルを間引き、
    「どの棒がどの AP か」が読めなくなる。**間引かせない**のが目的。
    """
    if count <= 12:
        return Pt(11)
    if count <= 24:
        return Pt(9)
    if count <= 32:
        return Pt(7)
    # 48 バケットを 5.5 インチ弱の作画領域に収めるには 5pt まで落とす必要がある
    # （6pt だと LibreOffice が 1 つおきに間引く。実際に PDF 化して確かめた値）
    return Pt(5)


def _reverse_categories(axis) -> None:
    """カテゴリ軸を反転して「先頭を上」にする。

    python-pptx には ``reverse_order`` に相当するプロパティが無いので
    ``c:scaling/c:orientation`` を直接書く。openpyxl 側（floorpeak / rrm の
    xlsx）で ``scaling.orientation = "maxMin"`` としているのと同じ意味。
    """
    scaling = axis._element.find(qn("c:scaling"))
    if scaling is None:  # pragma: no cover - python-pptx は必ず作る
        return
    orientation = scaling.find(qn("c:orientation"))
    if orientation is None:
        orientation = scaling.makeelement(qn("c:orientation"), {})
        scaling.insert(0, orientation)
    orientation.set("val", "maxMin")


def _apply_title(chart, title: str) -> None:
    chart.has_title = True
    frame = chart.chart_title.text_frame
    frame.text = title
    for run in frame.paragraphs[0].runs:
        run.font.size = Pt(12)
        run.font.bold = True


def _style(chart, *, category_count: int) -> None:
    """軸まわりの共通設定。**軸は消さない**（27 番の再発防止）。"""
    chart.font.size = Pt(9)
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = _label_size(category_count)
    _reverse_categories(category_axis)
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.size = Pt(9)


def add_stacked_bar(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    title: str,
    categories: Sequence[str],
    series: Sequence[tuple[str, Sequence[float]]],
    colors: Sequence[object] = (),
):
    """横棒の積み上げグラフ（系列ごとに色を指定する）。

    :param series: ``(系列名, 値の並び)`` の並び。値の数は ``categories`` と同数。
    """
    data = CategoryChartData()
    data.categories = list(categories)
    for name, values in series:
        data.add_series(name, list(values))

    frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, left, top, width, height, data)
    chart = frame.chart
    _apply_title(chart, title)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    _style(chart, category_count=len(categories))

    plot = chart.plots[0]
    plot.gap_width = GAP_WIDTH
    plot.overlap = 100  # 積み上げは 100 でないと横に並んでしまう
    for item, color in zip(plot.series, list(colors) + [None] * len(plot.series)):
        if color is None:
            continue
        item.format.fill.solid()
        item.format.fill.fore_color.rgb = _rgb(color)
    return chart


def add_bar(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    title: str,
    series_name: str,
    categories: Sequence[str],
    values: Sequence[float],
    point_colors: Sequence[object] | None = None,
    data_labels: bool = True,
):
    """単一系列の横棒グラフ。``point_colors`` を渡すと棒ごとに色を変える。

    単一系列では凡例に意味が無い（系列名しか出ない）ので凡例は出さない。
    色の意味づけが要る場合は呼び出し側でスライドに凡例を描くこと。
    """
    data = CategoryChartData()
    data.categories = list(categories)
    data.add_series(series_name, list(values))

    frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, data)
    chart = frame.chart
    _apply_title(chart, title)
    chart.has_legend = False
    _style(chart, category_count=len(categories))

    plot = chart.plots[0]
    plot.gap_width = GAP_WIDTH
    if data_labels:
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(9)
    if point_colors:
        for point, color in zip(plot.series[0].points, point_colors):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(color)
    return chart
