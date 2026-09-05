"""横断レポート（PPTX）の組み立てテスト。合成データのみを使う。

要点:

- 3 モジュールとも未選択はエラー（空のレポートを作らない）
- 選ばれたモジュールの章**だけ**が出る
- 章の順序は選んだ順ではなく **Hang AP → Floor Peak → RRM の固定順**
- 存在しない結果名はエラー
- 生成した pptx が python-pptx で読み戻せ、スライド数が期待どおり
- グラフの系列数・データ点数が保存済み結果と一致する
"""
from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path

import _repsynth as S
import pytest
from pptx import Presentation

from report import analysis, builder

GENERATED_AT = datetime(2026, 1, 4, 4, 4, 4, tzinfo=timezone.utc)


@pytest.fixture
def dirs(tmp_path) -> analysis.ResultsDirs:
    """3 モジュールぶんの保存済み結果を書いた保存先を返す。"""
    d = analysis.ResultsDirs.under(tmp_path)
    S.write_hangap(d.hangap)
    S.write_floorpeak(d.floorpeak)
    S.write_rrm(d.rrm)
    return d


def _build(dirs: analysis.ResultsDirs, **kwargs) -> analysis.ReportResult:
    params = analysis.ReportParams(**kwargs)
    return analysis.run_report(params, dirs, generated_at=GENERATED_AT)


def _charts(prs: Presentation) -> list:
    return [
        shape.chart
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_chart", False)
    ]


# ---------------------------------------------------------------------------
# 1. 3 モジュールすべて未選択
# ---------------------------------------------------------------------------


def test_no_selection_is_error(dirs):
    with pytest.raises(analysis.ParamError):
        _build(dirs)


def test_blank_names_count_as_unselected(dirs):
    """空文字は「選ばれていない」として扱う（空のレポートにしない）。"""
    with pytest.raises(analysis.ParamError):
        _build(dirs, hangap_result="", floorpeak_result="   ", rrm_result=None)


# ---------------------------------------------------------------------------
# 2. 1 モジュールだけ選んだ場合、そのモジュールのスライドだけが出る
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("field", "section"),
    [
        ("hangap_result", "hangap"),
        ("floorpeak_result", "floorpeak"),
        ("rrm_result", "rrm"),
    ],
)
def test_single_section_only(dirs, field, section):
    name = {
        "hangap_result": S.HANGAP_NAME,
        "floorpeak_result": S.FLOORPEAK_NAME,
        "rrm_result": S.RRM_NAME,
    }[field]
    result = _build(dirs, **{field: name})

    assert result.sections == [section]
    sections = {s.section for s in result.slides}
    assert sections == {"cover", section}
    # 他モジュールのスライドが混ざっていないこと
    for other in set(analysis.SECTION_ORDER) - {section}:
        assert other not in sections


def test_two_sections_only(dirs):
    result = _build(dirs, hangap_result=S.HANGAP_NAME, rrm_result=S.RRM_NAME)
    assert result.sections == ["hangap", "rrm"]
    assert "floorpeak" not in {s.section for s in result.slides}


# ---------------------------------------------------------------------------
# 3. 章立ての順序は固定（選んだ順序で変わらない）
# ---------------------------------------------------------------------------


def test_section_order_is_fixed(dirs):
    """パラメータの並びを変えても、章の順序は SECTION_ORDER のまま。"""
    forward = _build(
        dirs, hangap_result=S.HANGAP_NAME, floorpeak_result=S.FLOORPEAK_NAME,
        rrm_result=S.RRM_NAME,
    )
    backward = _build(
        dirs, rrm_result=S.RRM_NAME, floorpeak_result=S.FLOORPEAK_NAME,
        hangap_result=S.HANGAP_NAME,
    )
    assert forward.sections == ["hangap", "floorpeak", "rrm"]
    assert backward.sections == forward.sections
    assert [s.section for s in backward.slides] == [s.section for s in forward.slides]


def test_section_order_is_fixed_even_when_sources_are_shuffled(dirs):
    """``build_report`` に渡す順序も結果に影響しない。"""
    params = analysis.ReportParams(
        hangap_result=S.HANGAP_NAME, floorpeak_result=S.FLOORPEAK_NAME, rrm_result=S.RRM_NAME,
    )
    sources = analysis.load_sources(params, dirs)
    result = analysis.build_report(list(reversed(sources)), generated_at=GENERATED_AT)
    assert result.sections == ["hangap", "floorpeak", "rrm"]


# ---------------------------------------------------------------------------
# 4. 存在しない結果名
# ---------------------------------------------------------------------------


def test_missing_result_is_error(dirs):
    with pytest.raises(analysis.ReportError) as excinfo:
        _build(dirs, rrm_result="rrm_result_20991231_235959")
    assert "見つかりません" in str(excinfo.value)


def test_invalid_name_is_param_error(dirs):
    """パス区切りや別モジュールの接頭辞は名前の検証で弾く。"""
    for bad in ("../rrm_result_20260103_030303", "hangap_result_20260101_010101"):
        with pytest.raises(analysis.ParamError):
            _build(dirs, rrm_result=bad)


def test_missing_csv_is_error(tmp_path):
    """json だけあって csv が無い組はエラー（明細のないレポートを作らない）。"""
    dirs = analysis.ResultsDirs.under(tmp_path)
    S.write_rrm(dirs.rrm)
    (dirs.rrm / f"{S.RRM_NAME}.csv").unlink()
    with pytest.raises(analysis.ReportError) as excinfo:
        _build(dirs, rrm_result=S.RRM_NAME)
    assert "csv" in str(excinfo.value)


# ---------------------------------------------------------------------------
# 5. 生成した pptx が読み戻せ、スライド数が期待どおり
# ---------------------------------------------------------------------------


def test_written_pptx_slide_count(dirs, tmp_path):
    result = _build(
        dirs, hangap_result=S.HANGAP_NAME, floorpeak_result=S.FLOORPEAK_NAME,
        rrm_result=S.RRM_NAME,
    )
    out = analysis.write_pptx(tmp_path / "out" / "report.pptx", result)
    assert out.is_file()

    prs = Presentation(str(out))
    # 表紙 1 + Hang AP 3 + Floor Peak（サマリ + フロア 2 枚）+ RRM 3
    assert len(prs.slides) == 1 + 3 + (1 + 2) + 3
    assert len(prs.slides) == result.slide_count


def test_single_section_pptx_slide_count(dirs, tmp_path):
    result = _build(dirs, rrm_result=S.RRM_NAME)
    out = analysis.write_pptx(tmp_path / "rrm.pptx", result)
    prs = Presentation(str(out))
    assert len(prs.slides) == 1 + 3  # 表紙 + RRM 3 枚
    assert result.slide_count == len(prs.slides)


def test_output_name_uses_generated_at():
    assert analysis.output_name(GENERATED_AT) == "report_20260104_040404.pptx"


# ---------------------------------------------------------------------------
# 6. グラフの系列数・データ点数が保存済み結果と一致する
# ---------------------------------------------------------------------------


def test_rrm_chart_series_and_points_match_meta(dirs, tmp_path):
    result = _build(dirs, rrm_result=S.RRM_NAME)
    source = result.sources[0]
    hourly = source.meta["hourly"]
    classifications = source.meta["classifications"]

    prs = Presentation(str(analysis.write_pptx(tmp_path / "rrm.pptx", result)))
    charts = _charts(prs)
    assert len(charts) == 2  # 時間帯別（積み上げ）+ 分類別インパクト

    hourly_chart, impact_chart = charts
    series = list(hourly_chart.plots[0].series)
    assert [s.name for s in series] == classifications
    for item, name in zip(series, classifications):
        assert list(item.values) == [b[f"changes_{name}"] for b in hourly]
    assert len(list(hourly_chart.plots[0].categories)) == len(hourly)

    by_class = source.meta["by_classification"]
    impact_series = list(impact_chart.plots[0].series)
    assert len(impact_series) == 1
    assert list(impact_series[0].values) == [c["impact_total"] for c in by_class]


def test_rrm_chart_caps_buckets(tmp_path):
    """バケットが上限を超えたら **新しい側** だけを描き、注記に全件数を残す。"""
    dirs = analysis.ResultsDirs.under(tmp_path)
    total = builder.MAX_CHART_BUCKETS + 12
    S.write_rrm(dirs.rrm, hourly=S.rrm_hourly(total))
    result = _build(dirs, rrm_result=S.RRM_NAME)

    prs = Presentation(str(analysis.write_pptx(tmp_path / "rrm.pptx", result)))
    hourly_chart = _charts(prs)[0]
    assert len(list(hourly_chart.plots[0].categories)) == builder.MAX_CHART_BUCKETS

    hourly = result.sources[0].meta["hourly"]
    shown = builder.rrm_chart_buckets(result.sources[0].meta)
    assert shown == hourly[-builder.MAX_CHART_BUCKETS:]


def test_floorpeak_chart_points_match_rows(dirs, tmp_path):
    result = _build(dirs, floorpeak_result=S.FLOORPEAK_NAME)
    source = result.sources[0]
    floors = source.meta["floors"]

    prs = Presentation(str(analysis.write_pptx(tmp_path / "fp.pptx", result)))
    charts = _charts(prs)
    assert len(charts) == len(floors)  # フロアごとに 1 枚

    for chart, floor in zip(charts, floors):
        series = list(chart.plots[0].series)
        assert len(series) == 1
        rows = source.rows[source.rows["map_name"] == floor["map_name"]]
        rows = rows.sort_values("rank_in_floor").head(builder.FLOORPEAK_TOP_N)
        assert list(series[0].values) == [float(v) for v in rows["num_clients"]]
        assert list(chart.plots[0].categories) == list(rows["ap_name"])


def test_floorpeak_chart_slides_are_capped(tmp_path):
    """フロアが多すぎるときはグラフのスライド数を打ち切る（打ち切りは注記する）。"""
    dirs = analysis.ResultsDirs.under(tmp_path)
    rows = [
        {
            "ap_name": f"TEST-AP-{i:02d}", "mac": f"aabbccddee{i:02d}", "model": "AP45",
            "num_clients": 1, "status": "connected", "map_id": f"test-map-{i}",
            "map_name": f"TEST-FLOOR-{i:02d}", "x_m": "", "y_m": "", "rank_in_floor": 1,
        }
        for i in range(1, builder.MAX_FLOOR_CHARTS + 4)
    ]
    S.write_floorpeak(dirs.floorpeak, rows=rows)
    result = _build(dirs, floorpeak_result=S.FLOORPEAK_NAME)
    charts_ = [s for s in result.slides if s.kind == "chart"]
    assert len(charts_) == builder.MAX_FLOOR_CHARTS


# ---------------------------------------------------------------------------
# 明細の要約（並べ替えと切り出し。**集計はしない**）
# ---------------------------------------------------------------------------


def test_hangap_detail_is_top_n_by_zero_count(tmp_path):
    dirs = analysis.ResultsDirs.under(tmp_path)
    S.write_hangap(dirs.hangap, rows=S.hangap_rows(builder.DETAIL_ROWS + 5))
    result = _build(dirs, hangap_result=S.HANGAP_NAME)
    rows = result.sources[0].rows

    top = builder._top_rows(rows, builder.HANGAP_SORT_COLUMN, builder.DETAIL_ROWS)
    assert len(top) == builder.DETAIL_ROWS
    values = [int(v) for v in top[builder.HANGAP_SORT_COLUMN]]
    assert values == sorted(values, reverse=True)
    assert values[0] == int(rows[builder.HANGAP_SORT_COLUMN].max())


# ---------------------------------------------------------------------------
# 明細の要約（直前clients の多い順。連続ゼロ回数トップ15 の直後に追加した2枚目）
# ---------------------------------------------------------------------------


def test_hangap_prev_clients_slide_is_added_after_zero_count_slide(dirs):
    """連続ゼロ回数トップ15 の直後に「直前clients」スライドが 1 枚増える。"""
    result = _build(dirs, hangap_result=S.HANGAP_NAME)
    hangap_slides = [s for s in result.slides if s.section == "hangap"]
    assert len(hangap_slides) == 3  # サマリ + 連続ゼロ回数トップ15 + 直前clientsトップ15
    assert builder.HANGAP_SORT_COLUMN in hangap_slides[1].title
    assert builder.HANGAP_PREV_CLIENTS_SORT_COLUMN in hangap_slides[2].title


def test_hangap_prev_clients_is_top_n_sorted_desc_with_name_tiebreak(tmp_path):
    """直前clients の降順。同値は AP 名の昇順(タイブレーク基準が既存に無いため)。"""
    dirs = analysis.ResultsDirs.under(tmp_path)
    rows = [
        {**row, "直前clients": 5} for row in S.hangap_rows(builder.DETAIL_ROWS + 5)
    ]
    # 先頭 3 件を同値にして、タイブレークが AP 名の昇順になることを確認する
    for i in range(3):
        rows[i]["直前clients"] = 99
    S.write_hangap(dirs.hangap, rows=rows)
    result = _build(dirs, hangap_result=S.HANGAP_NAME)
    all_rows = result.sources[0].rows

    top = builder._top_rows(
        all_rows, builder.HANGAP_PREV_CLIENTS_SORT_COLUMN, builder.DETAIL_ROWS,
        tiebreak="ap_name",
    )
    assert len(top) == builder.DETAIL_ROWS
    values = [int(v) for v in top[builder.HANGAP_PREV_CLIENTS_SORT_COLUMN]]
    assert values == sorted(values, reverse=True)

    tied = top[top[builder.HANGAP_PREV_CLIENTS_SORT_COLUMN] == 99]
    assert len(tied) == 3
    names = list(tied["ap_name"])
    assert names == sorted(names)


def test_hangap_prev_clients_slide_shows_fewer_than_15_when_available(tmp_path):
    dirs = analysis.ResultsDirs.under(tmp_path)
    S.write_hangap(dirs.hangap, rows=S.hangap_rows(4))
    result = _build(dirs, hangap_result=S.HANGAP_NAME)
    rows = result.sources[0].rows

    top = builder._top_rows(
        rows, builder.HANGAP_PREV_CLIENTS_SORT_COLUMN, builder.DETAIL_ROWS,
        tiebreak="ap_name",
    )
    assert len(top) == 4


def test_hangap_zero_count_slide_is_unchanged_by_new_slide(tmp_path):
    """新スライド追加が既存の「連続ゼロ回数トップ15」の内容・順序に影響しないこと。"""
    dirs = analysis.ResultsDirs.under(tmp_path)
    S.write_hangap(dirs.hangap, rows=S.hangap_rows(builder.DETAIL_ROWS + 5))
    result = _build(dirs, hangap_result=S.HANGAP_NAME)
    rows = result.sources[0].rows

    top = builder._top_rows(rows, builder.HANGAP_SORT_COLUMN, builder.DETAIL_ROWS)
    assert len(top) == builder.DETAIL_ROWS
    values = [int(v) for v in top[builder.HANGAP_SORT_COLUMN]]
    assert values == sorted(values, reverse=True)
    assert values[0] == int(rows[builder.HANGAP_SORT_COLUMN].max())

    hangap_slides = [s for s in result.slides if s.section == "hangap"]
    assert builder.HANGAP_SORT_COLUMN in hangap_slides[1].title
    assert hangap_slides[1].kind == "table"


def test_cover_is_always_present(dirs):
    result = _build(dirs, rrm_result=S.RRM_NAME)
    assert result.slides[0].section == "cover"
    assert result.slides[0].kind == "cover"


def test_slide_titles_carry_section_labels(dirs):
    result = _build(
        dirs, hangap_result=S.HANGAP_NAME, floorpeak_result=S.FLOORPEAK_NAME,
        rrm_result=S.RRM_NAME,
    )
    for slide in result.slides[1:]:
        assert slide.title.startswith(analysis.SECTION_LABELS[slide.section])


def test_hangap_pptx_tables_content(tmp_path):
    """PPTX 実体を読み戻し、既存の連続ゼロ回数トップ15スライドが変わらず、
    直前clientsトップ15スライドに両方の列（直前clients / AP最大clients）が
    降順で入っていることを確認する。"""
    dirs = analysis.ResultsDirs.under(tmp_path)
    S.write_hangap(dirs.hangap, rows=S.hangap_rows(5))
    result = _build(dirs, hangap_result=S.HANGAP_NAME)
    out = analysis.write_pptx(tmp_path / "hangap.pptx", result)
    prs = Presentation(str(out))

    # スライド構成: 0=表紙, 1=サマリ, 2=連続ゼロ回数トップ15, 3=直前clientsトップ15
    zero_count_slide, prev_clients_slide = prs.slides[2], prs.slides[3]

    zero_count_table = next(s.table for s in zero_count_slide.shapes if s.has_table)
    zero_count_header = [zero_count_table.cell(0, c).text for c in range(len(zero_count_table.columns))]
    assert zero_count_header == list(builder.HANGAP_DETAIL_COLUMNS)

    prev_clients_table = next(s.table for s in prev_clients_slide.shapes if s.has_table)
    header = [prev_clients_table.cell(0, c).text for c in range(len(prev_clients_table.columns))]
    assert header == list(builder.HANGAP_PREV_CLIENTS_DETAIL_COLUMNS)
    assert "直前clients" in header
    assert "AP最大clients" in header

    prev_idx = header.index("直前clients")
    values = [
        int(prev_clients_table.cell(r, prev_idx).text) for r in range(1, len(prev_clients_table.rows))
    ]
    assert values == sorted(values, reverse=True)
